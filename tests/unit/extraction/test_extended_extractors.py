"""Tests for extended document extractors (PPTX, EML, XML, Images, HTML, XLSX, ODT, RTF, MSG)."""

from pathlib import Path

import pytest

from bloginator.extraction._extended_extractors import (
    extract_text_from_eml,
    extract_text_from_html,
    extract_text_from_image,
    extract_text_from_msg,
    extract_text_from_odt,
    extract_text_from_pptx,
    extract_text_from_rtf,
    extract_text_from_xlsx,
    extract_text_from_xml,
)


class TestExtractTextFromPptx:
    """Test PowerPoint PPTX extraction."""

    def test_extract_pptx_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent PPTX file raises error."""
        nonexistent = tmp_path / "nonexistent.pptx"

        with pytest.raises(FileNotFoundError):
            extract_text_from_pptx(nonexistent)

    @pytest.mark.skipif(
        not pytest.importorskip("pptx", reason="python-pptx not installed"),
        reason="python-pptx not installed",
    )
    def test_extract_pptx_basic(self, tmp_path: Path) -> None:
        """Test extracting text from a basic PPTX file."""
        from pptx import Presentation
        from pptx.util import Inches

        # Create a simple PPTX
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add a text box
        left = Inches(1)
        top = Inches(1)
        width = Inches(5)
        height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Hello from PowerPoint"

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        text = extract_text_from_pptx(pptx_path)

        assert "Hello from PowerPoint" in text
        assert "Slide 1" in text


class TestExtractTextFromEml:
    """Test Email (.eml) extraction."""

    def test_extract_eml_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent EML file raises error."""
        nonexistent = tmp_path / "nonexistent.eml"

        with pytest.raises(FileNotFoundError):
            extract_text_from_eml(nonexistent)

    def test_extract_eml_plain_text(self, tmp_path: Path) -> None:
        """Test extracting from plain text email."""
        eml_content = """From: sender@example.com
To: recipient@example.com
Subject: Test Email Subject
Date: Thu, 5 Dec 2024 10:00:00 +0000
Content-Type: text/plain; charset="utf-8"

This is the email body content.
It has multiple lines.
"""
        eml_path = tmp_path / "test.eml"
        eml_path.write_text(eml_content)

        text = extract_text_from_eml(eml_path)

        assert "Subject: Test Email Subject" in text
        assert "From: sender@example.com" in text
        assert "This is the email body content" in text

    def test_extract_eml_html(self, tmp_path: Path) -> None:
        """Test extracting from HTML email."""
        eml_content = """From: sender@example.com
To: recipient@example.com
Subject: HTML Email
Date: Thu, 5 Dec 2024 10:00:00 +0000
Content-Type: text/html; charset="utf-8"

<html><body><h1>Hello</h1><p>This is HTML content.</p></body></html>
"""
        eml_path = tmp_path / "test.eml"
        eml_path.write_text(eml_content)

        text = extract_text_from_eml(eml_path)

        assert "Hello" in text
        assert "HTML content" in text
        assert "<html>" not in text  # HTML should be stripped


class TestExtractTextFromXml:
    """Test XML extraction."""

    def test_extract_xml_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent XML file raises error."""
        nonexistent = tmp_path / "nonexistent.xml"

        with pytest.raises(FileNotFoundError):
            extract_text_from_xml(nonexistent)

    def test_extract_xml_basic(self, tmp_path: Path) -> None:
        """Test extracting text from basic XML."""
        xml_content = """<?xml version="1.0"?>
<document>
    <title>Test Document</title>
    <content>This is the main content.</content>
    <section>
        <paragraph>Nested paragraph text.</paragraph>
    </section>
</document>
"""
        xml_path = tmp_path / "test.xml"
        xml_path.write_text(xml_content)

        text = extract_text_from_xml(xml_path)

        assert "Test Document" in text
        assert "This is the main content" in text
        assert "Nested paragraph text" in text

    def test_extract_xml_invalid(self, tmp_path: Path) -> None:
        """Test extracting from invalid XML returns empty string (graceful degradation)."""
        xml_path = tmp_path / "invalid.xml"
        xml_path.write_text("This is not valid XML <unclosed>")

        # Should return empty string, not raise - graceful degradation
        result = extract_text_from_xml(xml_path)
        assert result == ""


class TestExtractTextFromImage:
    """Test image OCR extraction."""

    def test_extract_image_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent image file raises error."""
        nonexistent = tmp_path / "nonexistent.png"

        with pytest.raises(FileNotFoundError):
            extract_text_from_image(nonexistent)


class TestExtractTextFromHtml:
    """Test HTML extraction."""

    def test_extract_html_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent HTML file raises error."""
        nonexistent = tmp_path / "nonexistent.html"

        with pytest.raises(FileNotFoundError):
            extract_text_from_html(nonexistent)

    def test_extract_html_basic(self, tmp_path: Path) -> None:
        """Test extracting text from basic HTML."""
        html_content = """<!DOCTYPE html>
<html>
<head><title>Test Page</title></head>
<body>
    <h1>Main Heading</h1>
    <p>This is a paragraph with some text.</p>
    <div>Another section of content.</div>
</body>
</html>
"""
        html_path = tmp_path / "test.html"
        html_path.write_text(html_content)

        text = extract_text_from_html(html_path)

        assert "Main Heading" in text
        assert "This is a paragraph" in text
        assert "Another section" in text
        assert "<html>" not in text  # Tags should be stripped

    def test_extract_html_with_scripts(self, tmp_path: Path) -> None:
        """Test that script and style tags are removed."""
        html_content = """<html>
<head><script>var x = 1;</script><style>.foo { color: red; }</style></head>
<body><p>Visible content only</p></body>
</html>
"""
        html_path = tmp_path / "test.html"
        html_path.write_text(html_content)

        text = extract_text_from_html(html_path)

        assert "Visible content only" in text
        assert "var x" not in text
        assert "color: red" not in text


class TestExtractTextFromXlsx:
    """Test Excel XLSX extraction."""

    def test_extract_xlsx_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent XLSX file raises error."""
        nonexistent = tmp_path / "nonexistent.xlsx"

        with pytest.raises(FileNotFoundError):
            extract_text_from_xlsx(nonexistent)

    def test_extract_xlsx_basic(self, tmp_path: Path) -> None:
        """Test extracting text from basic XLSX file."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Header 1"
        ws["B1"] = "Header 2"
        ws["A2"] = "Data value 1"
        ws["B2"] = "Data value 2"

        xlsx_path = tmp_path / "test.xlsx"
        wb.save(xlsx_path)

        text = extract_text_from_xlsx(xlsx_path)

        assert "Header 1" in text
        assert "Header 2" in text
        assert "Data value 1" in text
        assert "Data value 2" in text

    def test_extract_xlsx_multiple_sheets(self, tmp_path: Path) -> None:
        """Test extracting from XLSX with multiple sheets."""
        from openpyxl import Workbook

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "First Sheet"
        ws1["A1"] = "First sheet content"

        ws2 = wb.create_sheet("Second Sheet")
        ws2["A1"] = "Second sheet content"

        xlsx_path = tmp_path / "multi.xlsx"
        wb.save(xlsx_path)

        text = extract_text_from_xlsx(xlsx_path)

        assert "First Sheet" in text
        assert "First sheet content" in text
        assert "Second Sheet" in text
        assert "Second sheet content" in text


class TestExtractTextFromRtf:
    """Test RTF extraction."""

    def test_extract_rtf_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent RTF file raises error."""
        nonexistent = tmp_path / "nonexistent.rtf"

        with pytest.raises(FileNotFoundError):
            extract_text_from_rtf(nonexistent)

    def test_extract_rtf_basic(self, tmp_path: Path) -> None:
        """Test extracting text from basic RTF file."""
        # Simple RTF content
        rtf_content = r"""{\rtf1\ansi
Hello World
\par
This is RTF content.
}"""
        rtf_path = tmp_path / "test.rtf"
        rtf_path.write_text(rtf_content)

        text = extract_text_from_rtf(rtf_path)

        # Should extract the text content
        assert "Hello" in text or "World" in text or "RTF content" in text


class TestExtractTextFromOdt:
    """Test OpenDocument Text extraction."""

    def test_extract_odt_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent ODT file raises error."""
        nonexistent = tmp_path / "nonexistent.odt"

        with pytest.raises(FileNotFoundError):
            extract_text_from_odt(nonexistent)

    def test_extract_odt_basic(self, tmp_path: Path) -> None:
        """Test extracting text from basic ODT file."""
        from odf import text as odf_text
        from odf.opendocument import OpenDocumentText

        doc = OpenDocumentText()
        para = odf_text.P(text="This is ODT paragraph content.")
        doc.text.addElement(para)

        para2 = odf_text.P(text="Another paragraph in the document.")
        doc.text.addElement(para2)

        odt_path = tmp_path / "test.odt"
        doc.save(str(odt_path))

        text = extract_text_from_odt(odt_path)

        assert "ODT paragraph content" in text
        assert "Another paragraph" in text


class TestExtractTextFromMsg:
    """Test Outlook MSG extraction."""

    def test_extract_msg_nonexistent(self, tmp_path: Path) -> None:
        """Test extracting from nonexistent MSG file raises error."""
        nonexistent = tmp_path / "nonexistent.msg"

        with pytest.raises(FileNotFoundError):
            extract_text_from_msg(nonexistent)

"""Extended document extractors for additional file formats.

Handles extraction from:
- PowerPoint (.ppt, .pptx)
- Email (.eml, .msg)
- XML (.xml)
- HTML (.html, .htm)
- Spreadsheets (.xlsx)
- Rich Text (.rtf)
- OpenDocument (.odt)
- Images (.png, .jpg, .jpeg, .webp) - via OCR when pytesseract is available
"""

import email
import shutil
import subprocess
from email import policy
from pathlib import Path
from xml.etree import ElementTree


def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract text content from PowerPoint PPTX file.

    Uses python-pptx to extract text from slides, including:
    - Slide titles
    - Text boxes
    - Notes

    Args:
        pptx_path: Path to PPTX file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If PPTX file does not exist
        ValueError: If file cannot be opened as PPTX
    """
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    try:
        from pptx import Presentation
    except ImportError as e:
        raise ValueError("python-pptx not installed") from e

    try:
        prs = Presentation(str(pptx_path))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- Slide {slide_num} ---"]

            # Extract text from shapes (text boxes, titles, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            # Extract notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Notes: {notes}]")

            if len(slide_texts) > 1:  # More than just the slide header
                text_parts.append("\n".join(slide_texts))

        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {e}") from e


def extract_text_from_ppt(ppt_path: Path) -> str:
    """Extract text from legacy .ppt file via LibreOffice conversion.

    Args:
        ppt_path: Path to legacy .ppt file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If .ppt file does not exist
        ValueError: If conversion fails
    """
    if not ppt_path.exists():
        raise FileNotFoundError(f".ppt file not found: {ppt_path}")

    # Import the LibreOffice helper from doc extractors
    from bloginator.extraction._doc_extractors import extract_doc_with_libreoffice

    try:
        return extract_doc_with_libreoffice(ppt_path)
    except Exception as e:
        raise ValueError(f"Failed to extract text from .ppt: {e}") from e


def extract_text_from_eml(eml_path: Path) -> str:
    """Extract text content from email (.eml) file.

    Parses the email and extracts:
    - Subject
    - From/To headers
    - Plain text body (or stripped HTML)

    Args:
        eml_path: Path to .eml file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If .eml file does not exist
        ValueError: If file cannot be parsed as email
    """
    if not eml_path.exists():
        raise FileNotFoundError(f"Email file not found: {eml_path}")

    try:
        with eml_path.open("rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)

        text_parts = []

        # Extract headers
        subject = msg.get("Subject", "")
        from_addr = msg.get("From", "")
        to_addr = msg.get("To", "")
        date = msg.get("Date", "")

        if subject:
            text_parts.append(f"Subject: {subject}")
        if from_addr:
            text_parts.append(f"From: {from_addr}")
        if to_addr:
            text_parts.append(f"To: {to_addr}")
        if date:
            text_parts.append(f"Date: {date}")

        text_parts.append("")  # Blank line before body

        # Extract body
        body = msg.get_body(preferencelist=("plain", "html"))
        if body:
            content = body.get_content()
            if isinstance(content, bytes):
                content = content.decode("utf-8", errors="replace")

            # If HTML, do basic stripping
            if body.get_content_type() == "text/html":
                from bloginator.extraction._doc_extractors import html_to_text

                content = html_to_text(content)

            text_parts.append(content)

        return "\n".join(text_parts).strip()
    except Exception as e:
        raise ValueError(f"Failed to extract text from email: {e}") from e


def extract_text_from_xml(xml_path: Path) -> str:
    """Extract text content from XML file.

    Recursively extracts all text content from XML elements.

    Args:
        xml_path: Path to XML file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If XML file does not exist
        ValueError: If file cannot be parsed as XML
    """
    if not xml_path.exists():
        raise FileNotFoundError(f"XML file not found: {xml_path}")

    try:
        tree = ElementTree.parse(xml_path)
        root = tree.getroot()

        # Recursively extract all text
        def get_text(element: ElementTree.Element) -> str:
            texts = []
            if element.text and element.text.strip():
                texts.append(element.text.strip())
            for child in element:
                texts.append(get_text(child))
                if child.tail and child.tail.strip():
                    texts.append(child.tail.strip())
            return " ".join(texts)

        return get_text(root)
    except ElementTree.ParseError as e:
        raise ValueError(f"Failed to parse XML: {e}") from e
    except Exception as e:
        raise ValueError(f"Failed to extract text from XML: {e}") from e


def extract_text_from_image(image_path: Path) -> str:
    """Extract text from image using OCR.

    Uses pytesseract for OCR when available.
    Supports: .png, .jpg, .jpeg, .webp

    Args:
        image_path: Path to image file

    Returns:
        Extracted text content (empty string if OCR unavailable)

    Raises:
        FileNotFoundError: If image file does not exist
        ValueError: If OCR fails
    """
    if not image_path.exists():
        raise FileNotFoundError(f"Image file not found: {image_path}")

    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        # OCR not available - return empty with warning
        return f"[OCR not available for {image_path.name}]"

    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text.strip()
    except pytesseract.TesseractNotFoundError:
        # Tesseract not installed
        return f"[Tesseract OCR not installed - cannot extract text from {image_path.name}]"
    except Exception as e:
        raise ValueError(f"Failed to extract text from image: {e}") from e


def extract_text_from_html(html_path: Path) -> str:
    """Extract text content from HTML file.

    Uses BeautifulSoup to parse HTML and extract readable text.

    Args:
        html_path: Path to HTML file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If HTML file does not exist
        ValueError: If file cannot be parsed as HTML
    """
    if not html_path.exists():
        raise FileNotFoundError(f"HTML file not found: {html_path}")

    try:
        from bs4 import BeautifulSoup

        content = html_path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(content, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style", "meta", "link"]):
            element.decompose()

        # Get text
        text = soup.get_text(separator="\n", strip=True)
        return text
    except ImportError:
        # Fall back to simple regex stripping
        from bloginator.extraction._doc_extractors import html_to_text

        content = html_path.read_text(encoding="utf-8", errors="replace")
        return html_to_text(content)
    except Exception as e:
        raise ValueError(f"Failed to extract text from HTML: {e}") from e


def extract_text_from_xlsx(xlsx_path: Path) -> str:
    """Extract text content from Excel XLSX file.

    Uses openpyxl to read cells and concatenate their text values.

    Args:
        xlsx_path: Path to XLSX file

    Returns:
        Extracted text content (each row on a line, cells tab-separated)

    Raises:
        FileNotFoundError: If XLSX file does not exist
        ValueError: If file cannot be opened as XLSX
    """
    if not xlsx_path.exists():
        raise FileNotFoundError(f"XLSX file not found: {xlsx_path}")

    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ValueError("openpyxl not installed") from e

    try:
        wb = load_workbook(xlsx_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(val.strip() for val in row_values):
                    text_parts.append("\t".join(row_values))

            text_parts.append("")  # Blank line between sheets

        wb.close()
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise ValueError(f"Failed to extract text from XLSX: {e}") from e


def extract_text_from_odt(odt_path: Path) -> str:
    """Extract text content from OpenDocument Text (.odt) file.

    Uses odfpy to extract paragraph text.

    Args:
        odt_path: Path to ODT file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If ODT file does not exist
        ValueError: If file cannot be opened as ODT
    """
    if not odt_path.exists():
        raise FileNotFoundError(f"ODT file not found: {odt_path}")

    try:
        from odf import text as odf_text
        from odf.opendocument import load as odf_load
    except ImportError as e:
        raise ValueError("odfpy not installed") from e

    try:
        doc = odf_load(str(odt_path))
        text_parts = []

        for para in doc.getElementsByType(odf_text.P):
            # Recursively get all text from paragraph
            para_text = "".join(
                node.data if hasattr(node, "data") else ""
                for node in para.childNodes
                if hasattr(node, "data")
            )
            # Also try direct text content
            if hasattr(para, "firstChild") and para.firstChild:
                direct = str(para)
                if direct.strip():
                    text_parts.append(direct.strip())
            elif para_text.strip():
                text_parts.append(para_text.strip())

        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from ODT: {e}") from e


def extract_text_from_rtf(rtf_path: Path) -> str:
    """Extract text content from Rich Text Format (.rtf) file.

    Tries unrtf CLI first (fast), then falls back to striprtf library.

    Args:
        rtf_path: Path to RTF file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If RTF file does not exist
        ValueError: If extraction fails
    """
    if not rtf_path.exists():
        raise FileNotFoundError(f"RTF file not found: {rtf_path}")

    # Try unrtf CLI first (installed via brew install textract)
    unrtf_path = shutil.which("unrtf")
    if unrtf_path:
        try:
            result = subprocess.run(
                ["unrtf", "--text", str(rtf_path.absolute())],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0 and result.stdout.strip():
                # unrtf outputs with some header lines, strip them
                lines = result.stdout.split("\n")
                # Skip header lines that start with ###
                content_lines = [ln for ln in lines if not ln.startswith("###")]
                return "\n".join(content_lines).strip()
        except (subprocess.TimeoutExpired, subprocess.SubprocessError):
            pass  # Fall through to other methods

    # Try striprtf library
    try:
        from striprtf.striprtf import rtf_to_text

        content = rtf_path.read_bytes()
        text = rtf_to_text(content.decode("utf-8", errors="replace"))
        return text.strip()
    except ImportError:
        pass

    # Last resort: basic regex stripping
    content = rtf_path.read_text(encoding="utf-8", errors="replace")
    import re

    # Remove RTF control words
    text = re.sub(r"\\[a-z]+\d*\s?", " ", content)
    # Remove braces
    text = re.sub(r"[{}]", "", text)
    # Clean up whitespace
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_text_from_msg(msg_path: Path) -> str:
    """Extract text content from Outlook MSG file.

    Uses extract-msg library to parse MSG files.

    Args:
        msg_path: Path to MSG file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If MSG file does not exist
        ValueError: If extraction fails
    """
    if not msg_path.exists():
        raise FileNotFoundError(f"MSG file not found: {msg_path}")

    try:
        import extract_msg
    except ImportError as e:
        raise ValueError("extract-msg not installed") from e

    try:
        msg = extract_msg.Message(str(msg_path))
        text_parts = []

        # Extract headers
        if msg.subject:
            text_parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            text_parts.append(f"From: {msg.sender}")
        if msg.to:
            text_parts.append(f"To: {msg.to}")
        if msg.date:
            text_parts.append(f"Date: {msg.date}")

        text_parts.append("")  # Blank line before body

        # Extract body (prefer plain text)
        if msg.body:
            text_parts.append(msg.body)
        elif msg.htmlBody:
            from bloginator.extraction._doc_extractors import html_to_text

            text_parts.append(html_to_text(msg.htmlBody))

        msg.close()
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise ValueError(f"Failed to extract text from MSG: {e}") from e

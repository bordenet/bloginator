"""Office document extractors for PowerPoint, Excel, and OpenDocument files."""

from pathlib import Path


def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract text content from PowerPoint PPTX file.

    Uses python-pptx to extract text from slides, including:
    - Slide titles
    - Text boxes
    - Notes

    Args:
        pptx_path: Path to PPTX file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If PPTX file does not exist
        ValueError: If file cannot be opened as PPTX
    """
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    try:
        from pptx import Presentation
    except ImportError as e:
        raise ValueError("python-pptx not installed") from e

    try:
        prs = Presentation(str(pptx_path))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- Slide {slide_num} ---"]

            # Extract text from shapes (text boxes, titles, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            # Extract notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Notes: {notes}]")

            if len(slide_texts) > 1:  # More than just the slide header
                text_parts.append("\n".join(slide_texts))

        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {e}") from e


def extract_text_from_ppt(ppt_path: Path) -> str:
    """Extract text from legacy .ppt file via LibreOffice conversion.

    Args:
        ppt_path: Path to legacy .ppt file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If .ppt file does not exist
        ValueError: If conversion fails
    """
    if not ppt_path.exists():
        raise FileNotFoundError(f".ppt file not found: {ppt_path}")

    # Import the LibreOffice helper from doc extractors
    from bloginator.extraction._doc_extractors import extract_doc_with_libreoffice

    try:
        return extract_doc_with_libreoffice(ppt_path)
    except Exception as e:
        raise ValueError(f"Failed to extract text from .ppt: {e}") from e


def extract_text_from_xlsx(xlsx_path: Path) -> str:
    """Extract text content from Excel XLSX file.

    Uses openpyxl to read cells and concatenate their text values.

    Args:
        xlsx_path: Path to XLSX file

    Returns:
        Extracted text content (each row on a line, cells tab-separated)

    Raises:
        FileNotFoundError: If XLSX file does not exist
        ValueError: If file cannot be opened as XLSX
    """
    if not xlsx_path.exists():
        raise FileNotFoundError(f"XLSX file not found: {xlsx_path}")

    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ValueError("openpyxl not installed") from e

    try:
        wb = load_workbook(xlsx_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(val.strip() for val in row_values):
                    text_parts.append("\t".join(row_values))

            text_parts.append("")  # Blank line between sheets

        wb.close()
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise ValueError(f"Failed to extract text from XLSX: {e}") from e


def extract_text_from_odt(odt_path: Path) -> str:
    """Extract text content from OpenDocument Text (.odt) file.

    Uses odfpy to extract paragraph text.

    Args:
        odt_path: Path to ODT file

    Returns:
        Extracted text content

    Raises:
        FileNotFoundError: If ODT file does not exist
        ValueError: If file cannot be opened as ODT
    """
    if not odt_path.exists():
        raise FileNotFoundError(f"ODT file not found: {odt_path}")

    try:
        from odf import text as odf_text
        from odf.opendocument import load as odf_load
    except ImportError as e:
        raise ValueError("odfpy not installed") from e

    try:
        doc = odf_load(str(odt_path))
        text_parts = []

        for para in doc.getElementsByType(odf_text.P):
            # Recursively get all text from paragraph
            para_text = "".join(
                node.data if hasattr(node, "data") else ""
                for node in para.childNodes
                if hasattr(node, "data")
            )
            # Also try direct text content
            if hasattr(para, "firstChild") and para.firstChild:
                direct = str(para)
                if direct.strip():
                    text_parts.append(direct.strip())
            elif para_text.strip():
                text_parts.append(para_text.strip())

        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from ODT: {e}") from e
